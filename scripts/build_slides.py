"""Sinh bộ slide thuyết trình (PPTX) tiếng Việt cho đề tài Traffic Sign Detection — YOLO vs DETR.

Bộ slide được thiết kế cho một buổi báo cáo giữa kỳ: vừa chặt chẽ về kỹ thuật, vừa dễ hiểu
cho người chưa biết gì về deep learning. Mỗi khái niệm khó đều đi kèm một ví dụ thực tế /
phép so sánh đời thường (hộp "Hãy hình dung" hoặc "Ví dụ thực tế"). Bộ slide chia làm 3 phần
cho 3 thành viên (Loi · Vinh · Tu), mỗi slide kèm "văn nói" (speaker notes) tiếng Việt để
người trình bày biết cần nói gì.

Các ảnh kết quả EDA / training thật trong `results/` được nhúng sẵn; nếu thiếu ảnh thì slide
vẫn dựng bình thường (chỉ bỏ qua ảnh đó).

Chạy:
    python -m scripts.build_slides
    # hoặc
    python scripts/build_slides.py

Đầu ra: reports/traffic_sign_slides.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

from src.utils.paths import REPO_ROOT, RESULTS_EDA, RESULTS_SAMPLES

# ── Vị trí đầu ra ────────────────────────────────────────────────────────────
OUT_DIR = REPO_ROOT / "reports"
OUT_PATH = OUT_DIR / "traffic_sign_slides.pptx"

# ── Bảng màu ─────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x2E, 0x6D, 0xB4)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
ORANGE = RGBColor(0xD9, 0x7A, 0x1E)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CALLOUT_BG = RGBColor(0xFD, 0xF4, 0xE3)   # vàng kem nhạt cho hộp ví dụ
CALLOUT_LINE = RGBColor(0xE3, 0xA8, 0x4E)

# Kích thước slide 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Ảnh mẫu annotate (chọn vài ảnh để minh hoạ dataset)
SAMPLE_IMAGES = [
    "road545_png.rf.02614a36cb366d36c8c23deca405fc02.jpg",
    "000430_jpg.rf.1008366e8e2805837605a83fc64bc470.jpg",
    "FisheyeCamera_1_00497_png.rf.241cc03642cdc931bcb50b1c4947595f.jpg",
]


def eda(name: str) -> Path:
    return RESULTS_EDA / name


def sample(name: str) -> Path:
    return RESULTS_SAMPLES / "annotated_train_samples" / name


# ── Helper dựng slide ────────────────────────────────────────────────────────
class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.blank = self.prs.slide_layouts[6]  # layout trống — tự kiểm soát mọi thứ

    # -- nội bộ ----------------------------------------------------------------
    def _new(self):
        return self.prs.slides.add_slide(self.blank)

    def _notes(self, slide, text: str) -> None:
        slide.notes_slide.notes_text_frame.text = text.strip()

    def _bg(self, slide, color: RGBColor) -> None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _band(self, slide, color: RGBColor, label: str | None = None) -> None:
        """Dải màu mỏng phía trên đầu slide (đánh dấu phần đang trình bày)."""
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.16))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        bar.shadow.inherit = False
        if label:
            box = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(7), Inches(0.4))
            tf = box.text_frame
            tf.text = label
            run = tf.paragraphs[0].runs[0]
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = color

    def _title(self, slide, text: str, top: float = 0.58) -> None:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(1.0))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = text
        p = tf.paragraphs[0]
        p.runs[0].font.size = Pt(28)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = NAVY

    def _callout(self, slide, label: str, text: str, top, width=Inches(12.3), left=Inches(0.5)) -> None:
        """Hộp 'Hãy hình dung / Ví dụ thực tế' — nền vàng kem, viền cam."""
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.15))
        box.fill.solid()
        box.fill.fore_color.rgb = CALLOUT_BG
        box.line.color.rgb = CALLOUT_LINE
        box.line.width = Pt(1.25)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"💡 {label}  "
        r1.font.size = Pt(14)
        r1.font.bold = True
        r1.font.color.rgb = ORANGE
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(14)
        r2.font.color.rgb = NAVY

    # -- slide bìa --------------------------------------------------------------
    def title_slide(self, title: str, subtitle: str, presenters: str, notes: str) -> None:
        slide = self._new()
        self._bg(slide, NAVY)
        box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.0))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = title
        r = tf.paragraphs[0].runs[0]
        r.font.size = Pt(40)
        r.font.bold = True
        r.font.color.rgb = WHITE

        sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.7), Inches(1.2))
        st = sub.text_frame
        st.word_wrap = True
        st.text = subtitle
        sr = st.paragraphs[0].runs[0]
        sr.font.size = Pt(20)
        sr.font.color.rgb = RGBColor(0xCF, 0xDD, 0xEE)

        pres = slide.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.0))
        pt = pres.text_frame
        pt.text = presenters
        pr = pt.paragraphs[0].runs[0]
        pr.font.size = Pt(16)
        pr.font.color.rgb = RGBColor(0xA9, 0xC2, 0xDE)
        self._notes(slide, notes)

    # -- slide section divider --------------------------------------------------
    def section_slide(self, part: str, title: str, agenda: list, color: RGBColor, notes: str) -> None:
        slide = self._new()
        self._bg(slide, color)
        box = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.8))
        tf = box.text_frame
        tf.text = part
        r = tf.paragraphs[0].runs[0]
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xDD, 0xEC, 0xF7)

        box2 = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.3))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        tf2.text = title
        r2 = tf2.paragraphs[0].runs[0]
        r2.font.size = Pt(34)
        r2.font.bold = True
        r2.font.color.rgb = WHITE

        if agenda:
            ab = slide.shapes.add_textbox(Inches(1.0), Inches(3.9), Inches(11.0), Inches(3.0))
            at = ab.text_frame
            at.word_wrap = True
            for i, item in enumerate(agenda):
                p = at.paragraphs[0] if i == 0 else at.add_paragraph()
                p.text = "→  " + item
                p.runs[0].font.size = Pt(16)
                p.runs[0].font.color.rgb = RGBColor(0xEA, 0xF2, 0xFB)
                p.space_after = Pt(6)
        self._notes(slide, notes)

    # -- slide bullet (tuỳ chọn ảnh phải + hộp ví dụ dưới) ---------------------
    def bullet_slide(
        self,
        title: str,
        bullets: list,
        notes: str,
        band_color: RGBColor,
        band_label: str,
        image: Path | None = None,
        callout: tuple | None = None,
    ) -> None:
        slide = self._new()
        self._band(slide, band_color, band_label)
        self._title(slide, title)

        has_img = image is not None and Path(image).exists()
        body_w = Inches(6.4) if has_img else Inches(12.3)
        body_h = Inches(4.2) if (callout and not has_img) else Inches(5.2)
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), body_w, body_h)
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(bullets):
            level = 0
            text = item
            if isinstance(item, tuple):
                text, level = item
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ("• " if level == 0 else "– ") + text
            p.level = level
            run = p.runs[0]
            run.font.size = Pt(17 if level == 0 else 15)
            run.font.color.rgb = NAVY if level == 0 else GREY
            p.space_after = Pt(5)

        if has_img:
            self._fit_image(slide, image, Inches(7.1), Inches(1.6), Inches(5.7), Inches(4.6))
            if callout:
                self._callout(slide, callout[0], callout[1], Inches(6.3), width=Inches(5.7), left=Inches(7.1))
        elif callout:
            self._callout(slide, callout[0], callout[1], Inches(6.05))
        self._notes(slide, notes)

    # -- slide 1+ ảnh lớn + caption + (tuỳ chọn) hộp ví dụ ----------------------
    def image_slide(
        self,
        title: str,
        images: list,
        caption: str,
        notes: str,
        band_color: RGBColor,
        band_label: str,
        callout: tuple | None = None,
    ) -> None:
        slide = self._new()
        self._band(slide, band_color, band_label)
        self._title(slide, title)
        imgs = [Path(p) for p in images if Path(p).exists()]
        img_bottom = Inches(1.65)
        img_h = Inches(3.9) if callout else Inches(4.6)
        if imgs:
            if len(imgs) == 1:
                self._fit_image(slide, imgs[0], Inches(2.4), img_bottom, Inches(8.5), img_h)
            else:
                n = len(imgs)
                gap = Inches(0.2)
                total_w = SLIDE_W - Inches(1.0) - gap * (n - 1)
                w = Emu(int(total_w / n))
                x = Inches(0.5)
                for p in imgs:
                    self._fit_image(slide, p, x, img_bottom, w, img_h)
                    x = Emu(int(x) + int(w) + int(gap))
        cap_top = Inches(5.7) if callout else Inches(6.45)
        if caption:
            box = slide.shapes.add_textbox(Inches(0.5), cap_top, Inches(12.3), Inches(0.7))
            tf = box.text_frame
            tf.word_wrap = True
            tf.text = caption
            r = tf.paragraphs[0].runs[0]
            r.font.size = Pt(13)
            r.font.italic = True
            r.font.color.rgb = GREY
        if callout:
            self._callout(slide, callout[0], callout[1], Inches(6.2))
        self._notes(slide, notes)

    # -- slide bảng + (tuỳ chọn) hộp ví dụ -------------------------------------
    def table_slide(
        self,
        title: str,
        headers: list,
        rows: list,
        notes: str,
        band_color: RGBColor,
        band_label: str,
        caption: str = "",
        callout: tuple | None = None,
    ) -> None:
        slide = self._new()
        self._band(slide, band_color, band_label)
        self._title(slide, title)
        n_rows = len(rows) + 1
        n_cols = len(headers)
        top = Inches(1.75)
        height = Inches(min(4.0, 0.46 * n_rows))
        table = slide.shapes.add_table(
            n_rows, n_cols, Inches(0.5), top, Inches(12.3), height
        ).table
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = h
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(15)
            para.runs[0].font.bold = True
            para.runs[0].font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = band_color
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = table.cell(r, c)
                cell.text = str(val)
                para = cell.text_frame.paragraphs[0]
                para.runs[0].font.size = Pt(13)
                para.runs[0].font.color.rgb = NAVY
        y = int(top) + int(height) + int(Inches(0.15))
        if caption:
            box = slide.shapes.add_textbox(Inches(0.5), Emu(y), Inches(12.3), Inches(0.5))
            tf = box.text_frame
            tf.word_wrap = True
            tf.text = caption
            r0 = tf.paragraphs[0].runs[0]
            r0.font.size = Pt(13)
            r0.font.italic = True
            r0.font.color.rgb = GREY
            y += int(Inches(0.55))
        if callout:
            self._callout(slide, callout[0], callout[1], Emu(min(y, int(Inches(6.2)))))
        self._notes(slide, notes)

    # -- nhúng ảnh giữ tỉ lệ, canh giữa trong khung ----------------------------
    def _fit_image(self, slide, path: Path, x, y, max_w, max_h) -> None:
        from PIL import Image

        try:
            with Image.open(path) as im:
                iw, ih = im.size
        except Exception:
            slide.shapes.add_picture(str(path), x, y, width=max_w)
            return
        ratio = min(int(max_w) / iw, int(max_h) / ih)
        w = int(iw * ratio)
        h = int(ih * ratio)
        cx = int(x) + (int(max_w) - w) // 2
        cy = int(y) + (int(max_h) - h) // 2
        slide.shapes.add_picture(str(path), Emu(cx), Emu(cy), width=Emu(w), height=Emu(h))

    def save(self) -> None:
        OUT_DIR.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(OUT_PATH))


# ── Nội dung deck ────────────────────────────────────────────────────────────
def build() -> None:
    d = Deck()

    # =================================================================
    # MỞ ĐẦU (chung)
    # =================================================================
    d.title_slide(
        "Nhận diện Biển báo Giao thông cho Xe Tự lái: YOLO vs DETR",
        "Báo cáo giữa kỳ — Môn Deep Learning · So sánh mô hình CNN một-giai-đoạn (YOLO) và "
        "Transformer (DETR) trên bài toán phát hiện biển báo",
        "Nhóm thực hiện: Loi · Vinh · Tu",
        "Xin chào thầy/cô và các bạn. Hôm nay nhóm em xin trình bày đề tài Nhận diện biển báo giao thông cho "
        "xe tự lái, so sánh hai họ mô hình là YOLO và DETR. Em xin phép trình bày theo hướng dễ hiểu nhất: "
        "với mỗi khái niệm kỹ thuật, nhóm sẽ kèm một ví dụ thực tế để cả những bạn chưa từng học deep "
        "learning cũng có thể theo dõi được. Bài gồm ba phần, do ba thành viên trình bày.",
    )

    d.bullet_slide(
        "Mở đầu: một tình huống quen thuộc",
        [
            "6 giờ tối, trời mưa, bạn lái xe trên một con đường lạ",
            "Phía trước có một biển báo — nhưng là biển STOP, hay biển giới hạn tốc độ? Đèn đỏ hay đèn xanh?",
            "Con người chỉ cần liếc mắt là biết. Nhưng MÁY TÍNH trên xe tự lái thì sao?",
            "Xe tự lái phải tự 'nhìn' qua camera và tự hiểu — chính xác, và trong tích tắc",
            "Đây chính là bài toán nhóm em nghiên cứu: dạy máy tính NHÌN và HIỂU biển báo",
        ],
        "Em muốn bắt đầu bằng một tình huống rất quen thuộc. Hãy tưởng tượng sáu giờ tối, trời mưa, bạn lái "
        "xe trên một con đường lạ. Phía trước có một biển báo. Là người, chúng ta chỉ cần liếc mắt là biết "
        "ngay đó là biển gì. Nhưng với một chiếc xe tự lái, nó phải tự nhìn qua camera và tự hiểu, vừa phải "
        "chính xác vừa phải nhanh trong tích tắc. Toàn bộ đề tài của nhóm xoay quanh việc dạy cho máy tính "
        "làm được đúng việc đó: nhìn và hiểu biển báo.",
        NAVY,
        "Mở đầu",
        callout=("Hãy hình dung:", "Nếu xe đọc nhầm biển 'giới hạn 30' thành 'giới hạn 80', hậu quả có thể "
                 "là một vụ tai nạn. Vì vậy bài toán này không chỉ là kỹ thuật — nó là AN TOÀN tính mạng."),
    )

    d.bullet_slide(
        "Câu hỏi nghiên cứu & Mục tiêu",
        [
            "Câu hỏi nghiên cứu: Làm sao nhận diện biển báo cho xe tự lái vừa CHÍNH XÁC vừa TIẾT KIỆM DỮ LIỆU, "
            "bằng YOLO và DETR?",
            ("'Chính xác' (robust): nhận đúng kể cả khi ảnh khó — mưa, tối, biển nhỏ", 1),
            ("'Tiết kiệm dữ liệu' (data-efficient): học tốt mà không cần quá nhiều ảnh gán nhãn", 1),
            "Mục tiêu giữa kỳ: dựng pipeline tái lập được + huấn luyện baseline cho CẢ HAI mô hình + so sánh",
            "So sánh trên 3 tiêu chí: độ chính xác (mAP), tốc độ (FPS), và kích thước mô hình (MB)",
        ],
        "Toàn bộ đề tài được dẫn dắt bởi một câu hỏi nghiên cứu: làm sao nhận diện biển báo vừa chính xác "
        "vừa tiết kiệm dữ liệu, bằng hai mô hình YOLO và DETR. Chính xác ở đây nghĩa là nhận đúng kể cả khi "
        "ảnh khó. Tiết kiệm dữ liệu nghĩa là học tốt mà không cần quá nhiều ảnh gán nhãn, vì gán nhãn rất "
        "tốn công và tiền. Ở giai đoạn giữa kỳ, mục tiêu của nhóm là dựng một pipeline tái lập được, huấn "
        "luyện baseline cho cả hai mô hình, và so sánh chúng trên ba tiêu chí: độ chính xác, tốc độ và kích thước.",
        NAVY,
        "Mở đầu",
    )

    d.bullet_slide(
        "Trước hết: 'Phát hiện đối tượng' là gì?",
        [
            "Phân loại ảnh (classification): chỉ trả lời 'trong ảnh CÓ con mèo' — một nhãn cho cả ảnh",
            "Phát hiện đối tượng (object detection): trả lời 'con mèo NẰM Ở ĐÂU' + 'đó là con gì'",
            ("→ Vừa khoanh vùng (vẽ khung), vừa dán nhãn cho TỪNG vật trong ảnh", 1),
            "Bài toán của nhóm là phát hiện đối tượng: tìm mọi biển báo trong ảnh và gọi tên từng cái",
            "Một ảnh có thể có 0, 1, hay nhiều biển báo — mô hình phải tìm ra hết",
        ],
        "Trước khi đi xa hơn, em xin làm rõ một khái niệm nền tảng cho các bạn chưa quen. Phân loại ảnh là "
        "khi máy chỉ trả lời trong ảnh có con mèo hay không, tức một nhãn cho cả ảnh. Còn phát hiện đối "
        "tượng thì khó hơn: nó phải trả lời con mèo nằm ở chỗ nào trong ảnh, và đó là con gì. Nghĩa là vừa "
        "khoanh vùng bằng một cái khung, vừa dán nhãn cho từng vật. Bài toán của nhóm thuộc loại thứ hai: "
        "tìm mọi biển báo trong một bức ảnh và gọi tên từng cái, dù ảnh có một hay nhiều biển.",
        NAVY,
        "Mở đầu",
        callout=("Ví dụ thực tế:", "Giống như khi bạn xem một bức ảnh lớp học và khoanh tròn từng bạn rồi "
                 "viết tên bên cạnh — chứ không chỉ nói 'đây là ảnh lớp học'."),
    )

    d.section_slide(
        "Bố cục buổi báo cáo",
        "Ba phần — ba người trình bày",
        [
            "Phần 1 (Loi): Đề tài, bài toán, bộ dữ liệu và phân tích khám phá dữ liệu (EDA)",
            "Phần 2 (Vinh): Mô hình YOLO — cách hoạt động, pipeline kỹ thuật, huấn luyện & kết quả",
            "Phần 3 (Tu): Mô hình DETR, so sánh YOLO vs DETR, bài học và hướng phát triển cuối kỳ",
        ],
        NAVY,
        "Bài trình bày chia làm ba phần. Phần một do Loi phụ trách về đề tài và dữ liệu. Phần hai do Vinh về "
        "mô hình YOLO. Phần ba do Tu về DETR, so sánh và hướng phát triển. Sau đây mời Loi bắt đầu phần một.",
    )

    # =================================================================
    # PHẦN 1 — LOI: Đề tài & Dữ liệu
    # =================================================================
    d.section_slide(
        "PHẦN 1 — Trình bày: Loi",
        "Đề tài & Bộ dữ liệu",
        [
            "Vì sao bài toán này quan trọng và khó",
            "Cách 'chấm điểm' một mô hình phát hiện (IoU, mAP, Precision/Recall)",
            "Bộ dữ liệu: 15 lớp biển báo và những gì EDA tiết lộ",
        ],
        BLUE,
        "Em là Loi, xin trình bày phần một. Em sẽ nói về lý do bài toán này quan trọng và khó, cách người ta "
        "chấm điểm một mô hình phát hiện, và những gì nhóm khám phá được từ bộ dữ liệu.",
    )

    d.bullet_slide(
        "Bối cảnh: Xe tự lái & vai trò của biển báo",
        [
            "Xe tự lái 'nhìn' thế giới qua camera, radar, lidar rồi tự ra quyết định lái",
            "Biển báo là tín hiệu điều khiển TRỰC TIẾP hành vi xe: dừng, giảm tốc, đi tiếp",
            "Đèn đỏ/xanh, biển STOP, biển giới hạn tốc độ — đọc sai là sai luật, là nguy hiểm",
            "Các hãng như Tesla, Waymo đều có module riêng chỉ để đọc biển báo & đèn tín hiệu",
            "Module này phải chạy LIÊN TỤC, real-time, trên máy tính nhỏ đặt trên xe",
        ],
        "Xe tự lái nhìn thế giới qua camera và các cảm biến rồi tự ra quyết định. Trong đó, biển báo là loại "
        "tín hiệu điều khiển trực tiếp hành vi của xe: khi nào dừng, khi nào giảm tốc, khi nào đi tiếp. Đọc "
        "sai một biển là vi phạm luật và gây nguy hiểm. Các hãng lớn như Tesla hay Waymo đều có riêng một "
        "module chỉ để đọc biển báo và đèn tín hiệu. Và module này phải chạy liên tục theo thời gian thực, "
        "trên một máy tính nhỏ gọn đặt ngay trên xe, chứ không phải siêu máy tính.",
        BLUE,
        "Phần 1 — Loi",
        callout=("Ví dụ thực tế:", "Năm 2020 từng có nghiên cứu cho thấy chỉ cần dán một mảnh băng keo nhỏ "
                 "lên biển tốc độ là khiến hệ thống đọc nhầm 35 thành 85. Đọc đúng biển báo là việc sống còn."),
    )

    d.bullet_slide(
        "Vì sao bài toán này KHÓ?",
        [
            "Biển báo thường rất NHỎ trong khung hình (ở xa, chiếm vài chục pixel)",
            "Điều kiện môi trường khắc nghiệt: mưa, sương, ngược nắng, ban đêm, đèn đường",
            "Biển bị che khuất một phần (cây, xe khác), bị mờ do xe đang chạy nhanh",
            "Nhiều biển TRÔNG GIỐNG NHAU: giới hạn 30, 80 chỉ khác mỗi con số",
            "Phải đúng VÀ nhanh: chần chừ vài trăm mili-giây cũng là quá muộn khi đang chạy",
        ],
        "Vì sao bài toán này khó? Thứ nhất, biển báo thường rất nhỏ trong khung hình vì ở xa, chỉ chiếm vài "
        "chục điểm ảnh. Thứ hai, điều kiện môi trường rất khắc nghiệt: mưa, sương, ngược nắng, ban đêm. Thứ "
        "ba, biển hay bị che khuất bởi cây hay xe khác, hoặc bị mờ vì xe đang chạy nhanh. Thứ tư, nhiều biển "
        "trông gần giống nhau, ví dụ giới hạn ba mươi và tám mươi chỉ khác con số. Và cuối cùng, mô hình "
        "phải vừa đúng vừa nhanh, vì chần chừ vài trăm mili-giây khi đang chạy đã là quá muộn.",
        BLUE,
        "Phần 1 — Loi",
        callout=("Hãy hình dung:", "Như việc bạn phải đọc một tấm biển nhỏ xíu, ướt mưa, ở cách 50 mét, "
                 "trong khi đang chạy bộ — và không được phép đọc sai."),
    )

    d.bullet_slide(
        "Cách 'chấm điểm': Bounding box & IoU",
        [
            "Bounding box = cái khung chữ nhật mô hình vẽ để khoanh vùng biển báo",
            "Làm sao biết khung mô hình vẽ có 'trúng' không? → so với khung đúng do người gán nhãn",
            "IoU (Intersection over Union) = phần GIAO chia phần HỢP của hai khung",
            ("IoU = 1.0: trùng khít hoàn hảo; IoU = 0: không chạm nhau", 1),
            ("Thường lấy ngưỡng IoU ≥ 0.5 thì coi là 'phát hiện đúng'", 1),
        ],
        "Bây giờ em nói về cách chấm điểm một mô hình. Khi mô hình tìm thấy biển báo, nó vẽ một cái khung "
        "chữ nhật gọi là bounding box. Làm sao biết khung đó vẽ có trúng không? Ta so nó với khung đúng do "
        "người gán nhãn, bằng một chỉ số tên là IoU, tức phần giao chia cho phần hợp của hai khung. Nếu hai "
        "khung trùng khít hoàn toàn thì IoU bằng một, nếu không chạm nhau thì bằng không. Thông thường nếu "
        "IoU từ 0.5 trở lên thì ta coi là phát hiện đúng vị trí.",
        BLUE,
        "Phần 1 — Loi",
        callout=("Hãy hình dung:", "Hai tờ giấy chồng lên nhau. Chồng càng khít, điểm càng cao. IoU đo đúng "
                 "mức độ chồng khít đó giữa khung-máy-vẽ và khung-đúng."),
    )

    d.bullet_slide(
        "Cách 'chấm điểm': Precision, Recall & mAP",
        [
            "Precision (độ chính xác): trong những gì máy BÁO là biển báo, bao nhiêu phần là ĐÚNG?",
            "Recall (độ bao phủ): trong tất cả biển báo THỰC SỰ có, máy tìm ra được bao nhiêu?",
            ("Precision cao mà recall thấp = ít báo nhầm nhưng bỏ sót nhiều", 1),
            ("Recall cao mà precision thấp = tìm được nhiều nhưng báo nhầm cũng nhiều", 1),
            "mAP (mean Average Precision): điểm tổng hợp cân bằng cả hai — CÀNG CAO CÀNG TỐT (tối đa 1.0)",
            "mAP@0.5 = chấm ở ngưỡng IoU 0.5; mAP@0.5:0.95 = chấm khắt khe hơn, trung bình nhiều ngưỡng",
        ],
        "Tiếp theo là ba chỉ số quan trọng. Precision trả lời: trong những gì máy báo là biển báo, bao "
        "nhiêu phần là đúng. Recall trả lời: trong tất cả biển báo thực sự có, máy tìm ra được bao nhiêu. "
        "Hai chỉ số này đánh đổi nhau. Để gộp lại thành một con số, người ta dùng mAP, càng cao càng tốt và "
        "tối đa là một. Trong báo cáo, các bạn sẽ thấy mAP ở ngưỡng 0.5, và mAP từ 0.5 đến 0.95 là cách "
        "chấm khắt khe hơn. Đây là những con số sẽ xuất hiện xuyên suốt phần kết quả.",
        BLUE,
        "Phần 1 — Loi",
        callout=("Ví dụ thực tế:", "Như bác sĩ tầm soát bệnh: Precision = trong số ca bị chẩn đoán dương "
                 "tính, bao nhiêu thật sự bệnh. Recall = trong số người thật sự bệnh, bác sĩ bắt được bao nhiêu."),
    )

    d.bullet_slide(
        "Hai 'nhân vật chính': YOLO và DETR",
        [
            "YOLO (You Only Look Once): mạng CNN, 'liếc' ảnh MỘT LẦN và đoán ngay tất cả khung + nhãn",
            ("Triết lý: nhanh, gọn, thực dụng — sinh ra để chạy real-time", 1),
            "DETR (DEtection TRansformer): dùng Transformer, 'suy nghĩ' về cả ảnh như một bài toán ghép cặp",
            ("Triết lý: thanh lịch, end-to-end, bỏ được nhiều bước thủ công — nhưng 'đói' dữ liệu", 1),
            "→ Hai triết lý đối lập trên cùng một bài toán: rất đáng để đặt lên bàn cân",
        ],
        "Cả bài có hai nhân vật chính. YOLO là một mạng tích chập, nó liếc qua ảnh chỉ một lần và đoán ngay "
        "tất cả các khung và nhãn cùng lúc, nên rất nhanh và thực dụng, sinh ra để chạy real-time. DETR thì "
        "dùng Transformer, nó suy nghĩ về toàn bộ ảnh như một bài toán ghép cặp thanh lịch, bỏ được nhiều "
        "bước thủ công, nhưng đổi lại rất đói dữ liệu. Hai triết lý đối lập này trên cùng một bài toán chính "
        "là lý do nhóm muốn đặt chúng lên bàn cân. Anh Vinh và anh Tu sẽ đi sâu vào từng mô hình ở phần sau.",
        BLUE,
        "Phần 1 — Loi",
    )

    d.bullet_slide(
        "Bộ dữ liệu nhóm sử dụng",
        [
            "Nguồn: pkdarabi/cardetection trên Kaggle — gốc từ nền tảng Roboflow (self-driving-cars v6)",
            "Lưu ý vui: tên là 'car detection' nhưng nhiệm vụ thực chất là PHÁT HIỆN BIỂN BÁO",
            "Ảnh đã chuẩn hoá về 416×416 điểm ảnh; nhãn theo định dạng chuẩn YOLOv8 (giấy phép CC BY 4.0)",
            "Tổng cộng: 4.969 ảnh và 6.012 biển báo được gán nhãn thủ công",
            "Giữ NGUYÊN 15 lớp gốc, không gộp/đổi nhãn → kết quả khách quan, dễ tái lập",
        ],
        "Bộ dữ liệu nhóm dùng là pkdarabi/cardetection trên Kaggle, gốc từ nền tảng Roboflow. Một lưu ý vui "
        "là tên nó nói car detection nhưng nhiệm vụ thật là phát hiện biển báo. Ảnh đã được chuẩn hoá về 416 "
        "trên 416 điểm ảnh, nhãn theo định dạng chuẩn của YOLOv8. Tổng cộng có gần năm nghìn ảnh với hơn sáu "
        "nghìn biển báo được con người gán nhãn thủ công. Nhóm cố ý giữ nguyên mười lăm lớp gốc, không gộp "
        "hay đổi nhãn, để kết quả khách quan và dễ tái lập cho người khác.",
        BLUE,
        "Phần 1 — Loi",
    )

    d.image_slide(
        "15 lớp biển báo & ảnh minh hoạ thật",
        [str(sample(s)) for s in SAMPLE_IMAGES],
        "15 lớp: Green Light, Red Light, Stop và 12 biển Speed Limit (10, 20, 30, 40, 50, 60, 70, 80, 90, "
        "100, 110, 120). Khung màu = nhãn do người gán, vẽ lại bằng src/data/visualize_annotations.py.",
        "Đây là vài ảnh thật từ bộ dữ liệu, có vẽ sẵn khung nhãn. Mười lăm lớp gồm đèn xanh, đèn đỏ, biển "
        "STOP, và mười hai biển giới hạn tốc độ từ mười đến một trăm hai mươi. Các bạn có thể thấy biển báo "
        "trong ảnh thực tế khá nhỏ, nhiều góc nhìn và nền phức tạp. Đây chính là minh hoạ cho những khó khăn "
        "mà em vừa nói ở các slide trước.",
        BLUE,
        "Phần 1 — Loi",
    )

    d.table_slide(
        "Chia dữ liệu: Train / Valid / Test",
        ["Tập", "Số ảnh", "Số biển báo", "Dùng để làm gì"],
        [
            ["Train (huấn luyện)", "3.530", "4.298", "Cho mô hình HỌC"],
            ["Valid (kiểm định)", "801", "944", "Theo dõi & tinh chỉnh trong lúc học"],
            ["Test (kiểm tra)", "638", "770", "Chấm điểm cuối — mô hình CHƯA từng thấy"],
            ["Tổng", "4.969", "6.012", "—"],
        ],
        "Dữ liệu được chia làm ba tập với vai trò khác nhau. Tập train để mô hình học. Tập valid để theo dõi "
        "và tinh chỉnh trong lúc học. Và tập test là tập mô hình chưa từng thấy, chỉ dùng để chấm điểm cuối "
        "cùng. Mọi con số kết quả mà nhóm báo cáo về sau đều đo trên tập test này, để đảm bảo công bằng và "
        "khách quan, giống như đề thi phải giữ kín tới lúc thi.",
        BLUE,
        "Phần 1 — Loi",
        caption="Số liệu đo trực tiếp từ bản dữ liệu thật — results/tables/dataset_summary.csv.",
        callout=("Ví dụ thực tế:", "Giống như ôn thi: bạn luyện trên đề cũ (train), tự kiểm tra bằng đề thử "
                 "(valid), nhưng điểm thật chỉ tính ở kỳ thi với đề chưa từng gặp (test)."),
    )

    d.bullet_slide(
        "EDA là gì & vì sao phải làm?",
        [
            "EDA (Exploratory Data Analysis) = phân tích khám phá dữ liệu TRƯỚC khi huấn luyện",
            "Mục đích: hiểu rõ dữ liệu — có bao nhiêu mỗi lớp, vật to hay nhỏ, nằm ở đâu, có lỗi gì không",
            "'Rác vào thì rác ra' — mô hình tốt đến đâu cũng không cứu được dữ liệu xấu",
            "EDA giúp dự đoán trước khó khăn và chọn đúng chiến lược huấn luyện",
            "Nhóm dùng các module trong src/eda/ để tự động sinh biểu đồ & bảng số liệu",
        ],
        "Trước khi huấn luyện, nhóm luôn làm một bước gọi là EDA, tức phân tích khám phá dữ liệu. Mục đích là "
        "hiểu thật rõ dữ liệu trong tay: mỗi lớp có bao nhiêu mẫu, vật thể to hay nhỏ, nằm ở đâu trong ảnh, "
        "và có lỗi gì không. Có một câu nói trong ngành là rác vào thì rác ra, nghĩa là mô hình giỏi đến đâu "
        "cũng không cứu được dữ liệu kém. EDA giúp nhóm lường trước khó khăn và chọn đúng chiến lược. Sau đây "
        "là vài phát hiện quan trọng nhất.",
        BLUE,
        "Phần 1 — Loi",
    )

    d.image_slide(
        "EDA — Phân bố lớp & sự mất cân bằng",
        [str(eda("class_distribution.png"))],
        "Lớp nhiều nhất 'Red Light' có 787 mẫu; lớp ít nhất 'Speed Limit 10' chỉ có 22 mẫu — chênh lệch 35.77 lần.",
        "Biểu đồ này cho thấy số lượng mẫu của từng lớp rất chênh lệch. Lớp đèn đỏ có tới bảy trăm tám mươi "
        "bảy mẫu, trong khi biển giới hạn mười chỉ có hai mươi hai mẫu, chênh nhau gần ba mươi sáu lần. Đây "
        "là vấn đề mất cân bằng lớp, và nó sẽ khiến mô hình học rất kém ở những lớp hiếm. Đây là một trong "
        "những phát hiện quan trọng nhất của EDA, và nó sẽ giải thích cho một số lỗi mà anh Vinh nêu ở phần sau.",
        BLUE,
        "Phần 1 — Loi",
        callout=("Hãy hình dung:", "Như một lớp học mà 787 buổi dạy về 'đèn đỏ' nhưng chỉ 22 buổi về 'biển "
                 "giới hạn 10' — học sinh tất nhiên sẽ giỏi cái thứ nhất và yếu cái thứ hai."),
    )

    d.image_slide(
        "EDA — Biển báo to hay nhỏ?",
        [str(eda("bbox_size_categories.png")), str(eda("bbox_wh.png"))],
        "Phân loại theo diện tích: Nhỏ 35.3% · Vừa 13.4% · Lớn 51.3%. Hơn 1/3 số biển báo thuộc nhóm 'nhỏ' — nhóm khó nhất.",
        "Hai biểu đồ này nói về kích thước biển báo. Hơn một phần ba số biển thuộc nhóm nhỏ, đây là nhóm khó "
        "phát hiện nhất vì chỉ chiếm vài chục điểm ảnh. Một nửa thuộc nhóm lớn, thường là biển ở gần camera. "
        "Biểu đồ bên phải cho thấy biển báo thường gần vuông. Thông tin này rất hữu ích: nó giải thích vì "
        "sao các biển ở xa hay bị bỏ sót, và gợi ý cho nhóm về kích thước ảnh đầu vào nên dùng khi huấn luyện.",
        BLUE,
        "Phần 1 — Loi",
        callout=("Hãy hình dung:", "Tìm một biển nhỏ vài chục pixel trong ảnh giống như tìm một con tem dán "
                 "đâu đó trên cả một bức tường — dễ sót hơn nhiều so với vật to."),
    )

    d.image_slide(
        "EDA — Biển báo thường nằm ở đâu?",
        [str(eda("object_center_heatmap.png")), str(eda("objects_per_image.png"))],
        "Bản đồ nhiệt: biển báo tập trung ở vùng trung tâm – phía trên khung hình. Đa số ảnh chỉ có 1–2 biển báo.",
        "Bản đồ nhiệt bên trái cho thấy tâm các biển báo tập trung ở vùng giữa và hơi phía trên của ảnh, "
        "đúng với góc nhìn camera gắn trên xe. Biểu đồ bên phải cho thấy đa số ảnh chỉ có một đến hai biển "
        "báo, hiếm khi dày đặc. Nghĩa là khung cảnh không quá rối, nhưng thử thách chính vẫn là vật thể nhỏ "
        "chứ không phải số lượng. Đây là những hiểu biết giúp nhóm đặt kỳ vọng đúng về độ khó.",
        BLUE,
        "Phần 1 — Loi",
    )

    d.bullet_slide(
        "Kiểm tra chất lượng dữ liệu (Data Quality)",
        [
            "Module src/data/validate_labels.py tự động soát: nhãn thiếu, nhãn rỗng, sai định dạng, ảnh hỏng, trùng lặp",
            "Phát hiện: 4 ảnh có nhãn RỖNG — file nhãn tồn tại nhưng không chứa biển báo nào",
            ("Đây là 'ảnh nền' (background) — ảnh không có biển báo, GIỮ LẠI sẽ giúp mô hình bớt báo nhầm", 1),
            "Không có ảnh hỏng, không có nhãn sai định dạng → dữ liệu nhìn chung SẠCH",
            "Kết quả lưu ở results/tables/data_quality_report.csv để bất kỳ ai cũng kiểm chứng được",
        ],
        "Bước cuối của phần dữ liệu là kiểm tra chất lượng. Nhóm có một module tự động soát các lỗi như nhãn "
        "thiếu, nhãn rỗng, sai định dạng, ảnh hỏng hay trùng lặp. Kết quả phát hiện bốn ảnh có nhãn rỗng, "
        "tức là ảnh không chứa biển báo nào. Thực ra đây là các ảnh nền, và việc giữ chúng lại còn có ích vì "
        "dạy mô hình biết khi nào không nên báo. Ngoài ra dữ liệu rất sạch, không có ảnh hỏng. Mọi kết quả "
        "đều được lưu thành file để người khác kiểm chứng.",
        BLUE,
        "Phần 1 — Loi",
        callout=("Ví dụ thực tế:", "Ảnh nền giống như câu hỏi 'bẫy' không có đáp án — dạy mô hình can đảm "
                 "nói 'ở đây không có biển báo' thay vì cố đoán bừa."),
    )

    d.bullet_slide(
        "Tóm tắt Phần 1",
        [
            "Bài toán: phát hiện biển báo cho xe tự lái — phải vừa CHÍNH XÁC vừa NHANH, vì liên quan an toàn",
            "Đã hiểu cách chấm điểm: IoU (vị trí), Precision/Recall (đúng/sót), mAP (điểm tổng hợp)",
            "Dữ liệu: 4.969 ảnh · 6.012 biển báo · 15 lớp · ảnh 416×416 · đã kiểm tra sạch",
            "Hai thách thức lớn nhất: mất cân bằng lớp 35.77× và rất nhiều biển báo NHỎ",
            "→ Mời Vinh trình bày mô hình đầu tiên: YOLO",
        ],
        "Tóm lại phần một: nhóm đã làm rõ bài toán phát hiện biển báo phải vừa chính xác vừa nhanh vì liên "
        "quan tới an toàn. Các bạn cũng đã nắm được cách chấm điểm qua IoU, precision, recall và mAP. Bộ dữ "
        "liệu gồm gần năm nghìn ảnh, mười lăm lớp, đã được kiểm tra sạch, với hai thách thức lớn nhất là mất "
        "cân bằng lớp và nhiều biển nhỏ. Em xin hết phần một, mời anh Vinh trình bày về mô hình YOLO.",
        BLUE,
        "Phần 1 — Loi",
    )

    # =================================================================
    # PHẦN 2 — VINH: YOLO
    # =================================================================
    d.section_slide(
        "PHẦN 2 — Trình bày: Vinh",
        "Mô hình YOLO — Cách hoạt động, Pipeline & Kết quả",
        [
            "YOLO 'nhìn một lần' nghĩa là gì, và vì sao nó nhanh",
            "Cách nhóm tổ chức code để tái lập được (pipeline + kiểm thử)",
            "Cấu hình huấn luyện và kết quả thật của YOLO",
        ],
        GREEN,
        "Em là Vinh, xin trình bày phần hai về YOLO. Em sẽ giải thích YOLO nhìn một lần nghĩa là gì, cách "
        "nhóm tổ chức code để tái lập, và kết quả thật mà YOLO đạt được.",
    )

    d.bullet_slide(
        "YOLO là gì? — Trực giác trước, kỹ thuật sau",
        [
            "YOLO = You Only Look Once = 'Chỉ nhìn một lần'",
            "Thay vì quét ảnh nhiều lần ở nhiều vùng, YOLO chia ảnh thành một LƯỚI ô vuông",
            "Mỗi ô tự hỏi: 'có biển báo ở đây không? nếu có thì khung bao nhiêu, lớp gì?'",
            "Tất cả ô trả lời CÙNG LÚC trong một lần chạy mạng → cực nhanh",
            "Nhóm dùng YOLOv8 (thư viện Ultralytics): bản hiện đại, anchor-free, dễ huấn luyện",
        ],
        "Em xin bắt đầu bằng trực giác. YOLO viết tắt của You Only Look Once, nghĩa là chỉ nhìn một lần. "
        "Thay vì quét ảnh nhiều lần ở nhiều vùng như các phương pháp cũ, YOLO chia ảnh thành một lưới các ô "
        "vuông. Mỗi ô tự hỏi: ở đây có biển báo không, nếu có thì khung bao nhiêu và là lớp gì. Điều hay là "
        "tất cả các ô trả lời cùng một lúc chỉ trong một lần chạy mạng, nên YOLO cực kỳ nhanh. Nhóm dùng "
        "phiên bản YOLOv8 của thư viện Ultralytics, là bản hiện đại và dễ huấn luyện.",
        GREEN,
        "Phần 2 — Vinh",
        callout=("Hãy hình dung:", "Như khi bạn liếc nhanh vào một bức ảnh và lập tức 'thấy' luôn mọi thứ — "
                 "chứ không soi từng góc một cách tuần tự. Đó là tinh thần 'nhìn một lần'."),
    )

    d.bullet_slide(
        "Một vấn đề nhỏ: NMS (lọc khung trùng)",
        [
            "Vì nhiều ô cạnh nhau cùng 'thấy' một biển báo → mô hình vẽ NHIỀU khung chồng lên một vật",
            "NMS (Non-Maximum Suppression) = bước dọn dẹp: giữ khung tự tin nhất, bỏ các khung trùng",
            "Đây là một bước thủ công GẮN THÊM sau mô hình YOLO",
            "Ghi nhớ điểm này: DETR ở phần sau sẽ KHÔNG cần bước NMS — đó là điểm khác biệt lớn",
        ],
        "YOLO có một vấn đề nhỏ. Vì nhiều ô cạnh nhau cùng nhìn thấy một biển báo, nên mô hình hay vẽ nhiều "
        "khung chồng lên cùng một vật. Để xử lý, người ta thêm một bước dọn dẹp tên là NMS, nghĩa là giữ lại "
        "khung tự tin nhất và bỏ các khung trùng còn lại. Đây là một bước thủ công gắn thêm sau mô hình. Em "
        "xin các bạn ghi nhớ điểm này, vì ở phần ba, DETR sẽ không cần bước NMS, và đó là một điểm khác biệt "
        "rất đáng chú ý giữa hai mô hình.",
        GREEN,
        "Phần 2 — Vinh",
        callout=("Hãy hình dung:", "Như 5 người cùng chỉ tay vào một chiếc xe và hô 'xe kìa!'. NMS là người "
                 "điều phối nói: 'rồi, chỉ tính một lần thôi'."),
    )

    d.bullet_slide(
        "Pipeline kỹ thuật của dự án",
        [
            "Toàn bộ logic nằm trong các module src/; notebook chỉ là lớp điều phối mỏng gọi lại chúng",
            ("src/data — đọc & kiểm tra dữ liệu, chuyển đổi định dạng", 1),
            ("src/eda — sinh biểu đồ phân tích (các hình ở Phần 1)", 1),
            ("src/training — huấn luyện YOLO & DETR; src/evaluation — chấm điểm & so sánh", 1),
            ("src/utils — quản lý đường dẫn, vẽ biểu đồ, cố định tính ngẫu nhiên (seed)", 1),
            "Mọi đường dẫn tập trung ở src/utils/paths.py → chạy y hệt nhau trên laptop và trên Colab",
        ],
        "Về mặt kỹ thuật, nhóm thiết kế theo nguyên tắc toàn bộ logic nằm trong các module src, còn notebook "
        "chỉ là một lớp điều phối mỏng gọi lại chúng. Code được chia rõ ràng thành nhóm xử lý dữ liệu, nhóm "
        "phân tích EDA, nhóm huấn luyện và đánh giá, và nhóm tiện ích. Một điểm nhóm khá tự hào là mọi đường "
        "dẫn được tập trung ở một chỗ, nên code chạy y hệt nhau dù trên laptop hay trên Google Colab, không "
        "phải sửa gì. Đây là nền tảng cho tính tái lập mà em sẽ nói ngay sau đây.",
        GREEN,
        "Phần 2 — Vinh",
    )

    d.bullet_slide(
        "Tái lập được (Reproducibility) — vì sao quan trọng?",
        [
            "Khoa học đòi hỏi: người khác chạy lại phải ra KẾT QUẢ GIỐNG mình",
            "Máy tính có nhiều yếu tố ngẫu nhiên (xáo trộn dữ liệu, khởi tạo trọng số…)",
            "Nhóm 'khoá' mọi ngẫu nhiên bằng seed_everything(42) → chạy lại luôn ra cùng kết quả",
            "Thư viện được ghim phiên bản (requirements.txt) để môi trường không đổi",
            "~43 bài kiểm thử tự động (pytest) chạy trong ~1 giây để bắt lỗi sớm",
        ],
        "Tái lập được là yêu cầu cốt lõi của khoa học: người khác chạy lại phải ra kết quả giống mình. Nhưng "
        "máy tính có nhiều yếu tố ngẫu nhiên, ví dụ thứ tự xáo trộn dữ liệu hay cách khởi tạo trọng số ban "
        "đầu. Nhóm khoá tất cả những ngẫu nhiên đó bằng một con số seed cố định là bốn mươi hai, nên chạy "
        "lại luôn ra cùng kết quả. Nhóm cũng ghim phiên bản các thư viện, và viết khoảng bốn mươi ba bài "
        "kiểm thử tự động chạy chỉ trong một giây để bắt lỗi sớm. Đây là điều một bài báo cáo nghiêm túc cần có.",
        GREEN,
        "Phần 2 — Vinh",
        callout=("Ví dụ thực tế:", "Như một công thức nấu ăn ghi rõ từng gram nguyên liệu: ai làm theo cũng "
                 "ra đúng món đó, không phụ thuộc may rủi."),
    )

    d.bullet_slide(
        "Chuẩn bị dữ liệu cho mô hình",
        [
            "YOLO đọc trực tiếp nhãn dạng text: mỗi dòng = (lớp, tâm-x, tâm-y, rộng, cao) đã chuẩn hoá 0–1",
            "DETR lại cần định dạng COCO → nhóm viết module convert_to_coco.py để chuyển đổi tự động",
            "Ảnh 416×416 từ dữ liệu gốc; YOLOv8 tự phóng về kích thước huấn luyện 640×640",
            "Augmentation (tăng cường dữ liệu): xoay, đổi sáng… để mô hình học đa dạng hơn",
            "NHƯNG có một phép augmentation nhóm CỐ TÌNH TẮT — slide sau sẽ giải thích vì sao",
        ],
        "Về chuẩn bị dữ liệu: YOLO đọc trực tiếp nhãn dạng text, mỗi dòng ghi lớp và toạ độ khung đã chuẩn "
        "hoá. DETR thì cần định dạng COCO nên nhóm viết một module chuyển đổi tự động. Ảnh gốc 416 được "
        "YOLOv8 tự phóng lên 640 khi huấn luyện. Thông thường người ta dùng augmentation, tức tăng cường dữ "
        "liệu bằng cách xoay hay đổi độ sáng để mô hình học đa dạng hơn. Nhưng có một phép augmentation mà "
        "nhóm cố tình tắt đi, và slide sau em sẽ giải thích vì sao nó lại quan trọng đến vậy.",
        GREEN,
        "Phần 2 — Vinh",
    )

    d.bullet_slide(
        "Vì sao TẮT lật ngang ảnh? (một quyết định quan trọng)",
        [
            "Lật ngang (horizontal flip) là augmentation phổ biến — thường giúp mô hình tổng quát tốt hơn",
            "NHƯNG biển báo nhiều khi CÓ HƯỚNG: mũi tên rẽ trái, rẽ phải, nhập làn…",
            "Lật ngang một biển 'rẽ phải' sẽ biến nó thành 'rẽ trái' → NHÃN BỊ SAI HOÀN TOÀN",
            "Với xe tự lái, học sai hướng = cực kỳ nguy hiểm",
            "→ Nhóm đặt fliplr = 0 (tắt). Đây là ví dụ KIẾN THỨC MIỀN (domain knowledge) quan trọng hơn mặc định",
        ],
        "Đây là một quyết định mà em đặc biệt muốn nhấn mạnh, vì nó cho thấy tư duy chứ không chỉ làm theo "
        "mặc định. Lật ngang ảnh là một augmentation rất phổ biến và thường có lợi. Nhưng biển báo nhiều khi "
        "có hướng, ví dụ mũi tên rẽ trái hay rẽ phải. Nếu lật ngang một biển rẽ phải, nó sẽ thành rẽ trái, "
        "tức là nhãn bị sai hoàn toàn. Với xe tự lái, dạy sai hướng là cực kỳ nguy hiểm. Vì vậy nhóm cố tình "
        "tắt phép lật này. Đây là một ví dụ điển hình cho việc hiểu biết về bài toán quan trọng hơn cài đặt mặc định.",
        GREEN,
        "Phần 2 — Vinh",
        callout=("Ví dụ thực tế:", "Lật ảnh chữ 'b' sẽ thành chữ 'd'. Với biển chỉ hướng cũng vậy — soi "
                 "gương là đổi nghĩa, không thể dùng để dạy máy."),
    )

    d.table_slide(
        "Cấu hình huấn luyện YOLO",
        ["Tham số", "Giá trị", "Ý nghĩa ngắn gọn"],
        [
            ["Mô hình", "YOLOv8n (nano)", "Bản nhỏ nhất → nhanh & nhẹ"],
            ["Khởi tạo", "pretrained", "Kế thừa kiến thức từ mô hình đã học sẵn"],
            ["Số epoch", "30", "Học trọn bộ dữ liệu 30 lượt"],
            ["Kích thước ảnh", "640 × 640", "Đủ lớn để thấy biển nhỏ"],
            ["Batch size", "16", "Số ảnh học mỗi bước"],
            ["fliplr", "0.0 (TẮT)", "Không lật — giữ đúng hướng biển"],
        ],
        "Đây là cấu hình huấn luyện. Nhóm chọn YOLOv8n bản nano nhỏ nhất để vừa nhanh vừa nhẹ, dùng trọng số "
        "pretrained nghĩa là kế thừa kiến thức từ một mô hình đã học sẵn trên kho ảnh khổng lồ. Nhóm huấn "
        "luyện ba mươi epoch, tức cho mô hình học trọn bộ dữ liệu ba mươi lượt, với ảnh đầu vào 640 đủ lớn "
        "để thấy biển nhỏ. Và như đã nói, phép lật ảnh được tắt hoàn toàn. Việc huấn luyện chạy trên GPU của "
        "Google Colab vì cần nhiều tính toán.",
        GREEN,
        "Phần 2 — Vinh",
        callout=("Hãy hình dung:", "'Pretrained' giống như tuyển một nhân viên đã có kinh nghiệm rồi đào tạo "
                 "thêm về biển báo — nhanh hơn nhiều so với dạy người chưa biết gì."),
    )

    d.image_slide(
        "Đường cong huấn luyện — đọc thế nào?",
        [str(eda("yolo_training_curves.png")), str(eda("training_curves.png"))],
        "Các đường 'loss' (sai số) đi XUỐNG = mô hình ngày càng ít sai. Các đường 'mAP' đi LÊN = ngày càng "
        "giỏi. Cả hai ổn định về cuối → hội tụ tốt, không overfitting rõ rệt.",
        "Em xin hướng dẫn cách đọc biểu đồ này cho các bạn chưa quen. Các đường loss, tức sai số, đi xuống "
        "nghĩa là mô hình ngày càng ít sai. Các đường mAP đi lên nghĩa là mô hình ngày càng giỏi. Điều nhóm "
        "mong muốn là cả hai đều ổn định về cuối, và đúng là như vậy. Điều này cho thấy mô hình đã hội tụ "
        "tốt và không bị overfitting, tức không phải học vẹt. Đây là cơ sở để tin vào kết quả ở slide tiếp theo.",
        GREEN,
        "Phần 2 — Vinh",
        callout=("Hãy hình dung:", "Như theo dõi điểm thi thử qua từng tuần ôn: điểm tăng dần rồi ổn định "
                 "nghĩa là bạn đã thật sự nắm bài, không phải ăn may."),
    )

    d.table_slide(
        "Kết quả YOLO trên tập Test",
        ["Chỉ số", "Giá trị", "Nghĩa là gì?"],
        [
            ["mAP@0.5", "0.954", "Gần như hoàn hảo ở ngưỡng dễ"],
            ["mAP@0.5:0.95", "0.806", "Vẫn rất tốt ở thang chấm khắt khe"],
            ["Precision", "0.908", "~91% khung báo ra là đúng"],
            ["Recall", "0.956", "Tìm được ~96% biển báo thật"],
            ["Tốc độ", "113 FPS (8.83 ms)", "Xử lý 113 ảnh/giây → dư sức real-time"],
            ["Kích thước", "6.25 MB", "Đủ nhẹ để nhúng lên thiết bị nhỏ"],
        ],
        "Và đây là kết quả của YOLO trên tập test. mAP ngưỡng 0.5 đạt 0.954, gần như hoàn hảo. Ngay cả ở "
        "thang chấm khắt khe hơn vẫn đạt 0.806. Precision khoảng chín mươi mốt phần trăm nghĩa là khung báo "
        "ra hầu hết đều đúng. Recall chín mươi sáu phần trăm nghĩa là tìm được gần hết biển báo thật. Đặc "
        "biệt nhất, mô hình xử lý được một trăm mười ba ảnh mỗi giây và chỉ nặng hơn sáu megabyte, tức dư "
        "sức chạy real-time trên một thiết bị nhỏ gắn trên xe. Đây là một kết quả baseline rất ấn tượng.",
        GREEN,
        "Phần 2 — Vinh",
        caption="Số liệu từ results/metrics/yolo_baseline.json — đo trên tập test.",
    )

    d.image_slide(
        "Ma trận nhầm lẫn — mô hình sai ở đâu?",
        [str(eda("yolo_confusion_matrix.png"))],
        "Đường chéo đậm = dự đoán đúng. Ô ngoài đường chéo = nhầm lẫn, chủ yếu giữa các biển tốc độ giống "
        "nhau và ở lớp hiếm 'Speed Limit 10' (chỉ 22 mẫu).",
        "Ma trận nhầm lẫn cho biết mô hình sai ở đâu. Cách đọc: đường chéo càng đậm thì càng nhiều dự đoán "
        "đúng, còn các ô nằm ngoài đường chéo là chỗ bị nhầm. Ở đây đường chéo rất đậm, chứng tỏ đa số đúng. "
        "Những nhầm lẫn còn lại chủ yếu xảy ra giữa các biển giới hạn tốc độ trông giống nhau, và ở lớp hiếm "
        "là biển giới hạn mười, vốn chỉ có hai mươi hai mẫu. Đây chính là hệ quả trực tiếp của vấn đề mất "
        "cân bằng lớp mà anh Loi đã nêu ở phần một.",
        GREEN,
        "Phần 2 — Vinh",
    )

    d.bullet_slide(
        "Tóm tắt Phần 2",
        [
            "YOLO 'nhìn một lần' nên rất nhanh; cần thêm bước NMS để lọc khung trùng",
            "Pipeline tái lập được: seed cố định, thư viện ghim phiên bản, ~43 bài kiểm thử",
            "Quyết định có chủ đích: TẮT lật ảnh để không phá hỏng biển có hướng",
            "Kết quả baseline rất mạnh: mAP@0.5 = 0.954 · 113 FPS · chỉ 6.25 MB",
            "Lỗi còn lại nằm ở lớp hiếm & các biển tốc độ giống nhau → mở đường cho phần cải thiện",
            "→ Mời Tu trình bày DETR và đặt hai mô hình lên bàn cân",
        ],
        "Tóm lại phần hai: YOLO nhìn một lần nên rất nhanh, nhưng cần thêm bước NMS. Nhóm đã xây pipeline tái "
        "lập được với seed cố định và kiểm thử đầy đủ, và có những quyết định kỹ thuật có chủ đích như tắt "
        "lật ảnh. Kết quả baseline của YOLO rất mạnh, vừa chính xác vừa nhanh và nhẹ. Các lỗi còn lại nằm ở "
        "lớp hiếm, mở đường cho phần cải thiện sau này. Em xin hết phần hai, mời anh Tu trình bày về DETR và "
        "so sánh hai mô hình.",
        GREEN,
        "Phần 2 — Vinh",
    )

    # =================================================================
    # PHẦN 3 — TU: DETR, So sánh & Hướng phát triển
    # =================================================================
    d.section_slide(
        "PHẦN 3 — Trình bày: Tu",
        "DETR, So sánh & Hướng phát triển",
        [
            "DETR — cách Transformer 'ghép cặp' để phát hiện vật thể",
            "Kết quả DETR & so sánh thẳng với YOLO",
            "Vì sao DETR thấp, bài học, và kế hoạch cho đồ án cuối kỳ",
        ],
        ORANGE,
        "Em là Tu, xin trình bày phần ba. Em sẽ giải thích DETR hoạt động ra sao, so sánh nó với YOLO, phân "
        "tích vì sao kết quả như vậy, và trình bày kế hoạch cho đồ án cuối kỳ.",
    )

    d.bullet_slide(
        "DETR là gì? — Trực giác",
        [
            "DETR = DEtection TRansformer (Facebook AI, 2020)",
            "Dùng Transformer — kiến trúc đứng sau ChatGPT — vốn giỏi 'nhìn toàn cảnh' và liên hệ các phần",
            "Ý tưởng: đưa cho mô hình một số 'phiếu hỏi' (object query) cố định, vd 100 phiếu",
            "Mỗi phiếu đi 'tìm' một vật thể trong ảnh; cuối cùng GHÉP mỗi phiếu với một biển báo thật",
            "Không cần lưới, không cần anchor, KHÔNG CẦN NMS → gọn gàng, end-to-end",
        ],
        "DETR viết tắt của Detection Transformer, do Facebook giới thiệu năm 2020. Nó dùng kiến trúc "
        "Transformer, chính là kiến trúc đứng sau ChatGPT, vốn rất giỏi nhìn toàn cảnh và liên hệ các phần "
        "với nhau. Ý tưởng của DETR là đưa cho mô hình một số phiếu hỏi cố định, ví dụ một trăm phiếu. Mỗi "
        "phiếu đi tìm một vật thể trong ảnh, và cuối cùng mô hình ghép mỗi phiếu với một biển báo thật. Nhờ "
        "cách làm này, DETR không cần lưới, không cần anchor, và đặc biệt không cần bước NMS như YOLO, nên "
        "rất gọn gàng và end-to-end.",
        ORANGE,
        "Phần 3 — Tu",
        callout=("Hãy hình dung:", "Như một quản lý phát 100 nhân viên đi khảo sát, rồi ghép mỗi nhân viên "
                 "với đúng một khách hàng — không ai bị trùng, không cần người dọn dẹp về sau (không NMS)."),
    )

    d.bullet_slide(
        "DETR ghép cặp thế nào? (bipartite matching)",
        [
            "Mô hình đưa ra 100 dự đoán; ảnh thực tế chỉ có vài biển báo thật",
            "Cần ghép: mỗi biển báo thật ↔ đúng MỘT dự đoán phù hợp nhất (1–1)",
            "Dùng thuật toán 'ghép cặp tối ưu' (Hungarian / bipartite matching) để chọn cách ghép tốt nhất",
            "Phần dự đoán không ghép được → bị gán nhãn 'không có vật' (no object)",
            "Cơ chế này thanh lịch nhưng cần NHIỀU vòng học để các phiếu 'phân vai' ổn định",
        ],
        "Vậy DETR ghép cặp thế nào? Mô hình luôn đưa ra một trăm dự đoán, nhưng ảnh thật chỉ có vài biển báo. "
        "Nên cần ghép mỗi biển thật với đúng một dự đoán phù hợp nhất, theo kiểu một đối một. Để làm việc "
        "này tối ưu, DETR dùng một thuật toán ghép cặp kinh điển tên là Hungarian matching. Những dự đoán "
        "thừa không ghép được sẽ bị gán là không có vật. Cơ chế này rất thanh lịch về mặt toán học, nhưng "
        "điểm yếu là cần rất nhiều vòng học thì các phiếu hỏi mới phân vai ổn định được. Đây là mấu chốt cho "
        "kết quả mà em sắp trình bày.",
        ORANGE,
        "Phần 3 — Tu",
    )

    d.table_slide(
        "Cấu hình huấn luyện DETR (baseline)",
        ["Tham số", "Giá trị", "Ghi chú"],
        [
            ["Mô hình", "detr-resnet-50", "Fine-tune từ bản pretrained"],
            ["Số epoch", "10", "RẤT ÍT (bài gốc dùng ~300–500)"],
            ["Batch size", "2", "Nhỏ — do giới hạn bộ nhớ GPU"],
            ["Learning rate", "1e-5 (backbone 1e-6)", "Học rất chậm, thận trọng"],
            ["Định dạng dữ liệu", "COCO JSON", "Chuyển từ nhãn YOLO"],
            ["Phần cứng", "GPU Colab", "Tài nguyên miễn phí, có hạn"],
        ],
        "Đây là cấu hình DETR. Nhóm fine-tune từ bản detr-resnet-50 có sẵn. Nhưng em xin lưu ý ngay hai con "
        "số: nhóm chỉ huấn luyện được mười epoch, trong khi bài báo gốc dùng tới ba đến năm trăm epoch; và "
        "batch chỉ bằng hai vì giới hạn bộ nhớ GPU miễn phí của Colab. Đây là một cấu hình rất khiêm tốn so "
        "với những gì DETR thực sự cần. Em nói trước điều này để các bạn hiểu đúng bối cảnh khi nhìn kết quả "
        "ở slide tiếp theo.",
        ORANGE,
        "Phần 3 — Tu",
        callout=("Hãy hình dung:", "DETR như một học sinh thông minh nhưng cần học rất nhiều buổi. Cho nó "
                 "học 10 buổi rồi đi thi thì điểm thấp là điều dễ hiểu."),
    )

    d.table_slide(
        "Kết quả DETR trên tập Test",
        ["Chỉ số", "Giá trị", "Nghĩa là gì?"],
        [
            ["mAP@0.5", "0.122", "Thấp — mô hình chưa hội tụ"],
            ["mAP@0.5:0.95", "0.099", "Thấp ở thang khắt khe"],
            ["Tốc độ", "15.7 FPS (63.83 ms)", "Chậm hơn YOLO ~7 lần"],
            ["Kích thước", "166 MB", "Nặng hơn YOLO ~27 lần"],
        ],
        "Và đây là kết quả DETR. mAP ngưỡng 0.5 chỉ đạt 0.122, thấp hơn nhiều so với YOLO. Tốc độ cũng chậm "
        "hơn khoảng bảy lần, và mô hình nặng tới một trăm sáu mươi sáu megabyte, gấp gần ba mươi lần YOLO. "
        "Nhóm hoàn toàn không bất ngờ với con số này. Ở slide sau em sẽ phân tích kỹ vì sao DETR lại thấp "
        "như vậy, và quan trọng là vì sao điều đó KHÔNG có nghĩa DETR là một mô hình tồi.",
        ORANGE,
        "Phần 3 — Tu",
        caption="Số liệu từ results/metrics/detr_baseline.json — đo trên tập test.",
    )

    d.table_slide(
        "Đặt lên bàn cân: YOLO vs DETR",
        ["Tiêu chí", "YOLO", "DETR", "Thắng"],
        [
            ["Độ chính xác (mAP@0.5)", "0.954", "0.122", "YOLO"],
            ["Chính xác khắt khe (mAP@0.5:0.95)", "0.806", "0.099", "YOLO"],
            ["Tốc độ (FPS)", "113", "15.7", "YOLO"],
            ["Độ trễ (ms/ảnh)", "8.83", "63.83", "YOLO"],
            ["Kích thước (MB)", "6.25", "166", "YOLO"],
        ],
        "Đây là bảng so sánh trực tiếp. Ở giai đoạn baseline này, phải thừa nhận YOLO thắng trên cả ba tiêu "
        "chí: chính xác hơn nhiều lần, nhanh hơn khoảng bảy lần, và nhẹ hơn hơn hai mươi lần. Tuy nhiên em "
        "muốn nhấn mạnh đây là so sánh ở điều kiện CHƯA cân bằng: DETR chưa được huấn luyện đủ. Vì vậy bảng "
        "này nói lên 'mô hình nào thực dụng hơn ngay bây giờ', chứ chưa phải 'mô hình nào tốt hơn về bản chất'.",
        ORANGE,
        "Phần 3 — Tu",
        caption="Số liệu tổng hợp từ results/tables/comparison.csv.",
    )

    d.bullet_slide(
        "Vì sao DETR thấp ở baseline? (phân tích trung thực)",
        [
            "DETR nổi tiếng 'đói dữ liệu' — bài báo gốc huấn luyện ~300–500 epoch; nhóm chỉ chạy được 10",
            "Batch size = 2 quá nhỏ → tín hiệu học bị 'nhiễu', khó ổn định",
            "Cơ chế ghép cặp cần nhiều vòng để 100 phiếu hỏi 'phân vai' rõ ràng — 10 epoch là chưa đủ",
            "~5.000 ảnh là NHỎ với DETR (vốn thiết kế cho COCO ~118.000 ảnh)",
            "→ Kết quả thấp phản ánh ĐIỀU KIỆN huấn luyện, KHÔNG phải giới hạn bản chất của DETR",
        ],
        "Em xin phân tích trung thực vì sao DETR thấp. Lý do gốc là DETR vốn đói dữ liệu và cần lịch huấn "
        "luyện rất dài. Bài báo gốc train tới ba đến năm trăm epoch, còn nhóm chỉ chạy được mười. Batch chỉ "
        "bằng hai khiến tín hiệu học bị nhiễu và khó ổn định. Cơ chế ghép cặp cũng cần nhiều vòng để các "
        "phiếu hỏi phân vai, mà mười epoch là chưa đủ. Thêm nữa, năm nghìn ảnh là quá nhỏ so với quy mô DETR "
        "được thiết kế. Vì vậy kết luận quan trọng là: con số thấp phản ánh điều kiện huấn luyện hạn chế, "
        "chứ không phải bản chất của DETR.",
        ORANGE,
        "Phần 3 — Tu",
        callout=("Ví dụ thực tế:", "Như đánh giá một vận động viên marathon qua 1 km đầu — họ cần đường dài "
                 "để bứt tốc. DETR cũng cần 'đường dài' (nhiều dữ liệu & epoch) mới phát huy."),
    )

    d.bullet_slide(
        "Bài học rút ra ở giai đoạn baseline",
        [
            "Với dữ liệu nhỏ & tài nguyên hạn chế, YOLO là lựa chọn thực dụng và hiệu quả vượt trội",
            "DETR thanh lịch về ý tưởng nhưng 'tốn kém': cần GPU mạnh, batch lớn, nhiều epoch",
            "Muốn so sánh CÔNG BẰNG, phải cấp NGÂN SÁCH huấn luyện tương đương cho cả hai",
            "Chất lượng & độ cân bằng dữ liệu ảnh hưởng trực tiếp tới lớp hiếm (vd Speed Limit 10)",
            "Pipeline tái lập giúp lặp lại thí nghiệm nhanh khi có thêm tài nguyên",
        ],
        "Từ giai đoạn này, nhóm rút ra mấy bài học. Thứ nhất, với dữ liệu nhỏ và tài nguyên hạn chế, YOLO là "
        "lựa chọn thực dụng và hiệu quả vượt trội. Thứ hai, DETR thanh lịch về ý tưởng nhưng tốn kém tài "
        "nguyên. Thứ ba, muốn so sánh công bằng thì phải cấp ngân sách huấn luyện tương đương cho cả hai, "
        "đây là điều nhóm sẽ làm ở cuối kỳ. Và thứ tư, chất lượng cùng độ cân bằng của dữ liệu ảnh hưởng "
        "trực tiếp tới các lớp hiếm. May mắn là pipeline tái lập giúp nhóm lặp lại thí nghiệm nhanh khi có "
        "thêm tài nguyên.",
        ORANGE,
        "Phần 3 — Tu",
    )

    d.bullet_slide(
        "Hạn chế hiện tại (nhìn thẳng vào điểm yếu)",
        [
            "Mất cân bằng lớp 35.77× → lớp hiếm như 'Speed Limit 10' học rất kém",
            "Nhiều biển báo nhỏ (35%) dễ bị bỏ sót",
            "DETR chưa hội tụ → so sánh hai mô hình chưa thật sự công bằng",
            "Chưa kiểm tra độ bền (robustness) với mưa, ban đêm, ảnh mờ, biển bị che",
            "Mới đánh giá trên tập test có sẵn, chưa thử trên video / tình huống thực tế ngoài phân phối",
        ],
        "Nhóm cũng xin nhìn thẳng vào những điểm yếu hiện tại. Mất cân bằng lớp khiến các lớp hiếm học rất "
        "kém. Nhiều biển nhỏ dễ bị bỏ sót. DETR chưa hội tụ nên so sánh chưa thật công bằng. Nhóm cũng chưa "
        "kiểm tra độ bền của mô hình trong điều kiện khó như mưa, ban đêm hay biển bị che. Và mới chỉ đánh "
        "giá trên tập test có sẵn, chưa thử trên video hay tình huống thực tế. Chính những hạn chế này định "
        "hình kế hoạch cho đồ án cuối kỳ mà em trình bày ngay sau đây.",
        ORANGE,
        "Phần 3 — Tu",
    )

    d.bullet_slide(
        "Hướng phát triển (Đồ án cuối kỳ) — 1/2: Mô hình & Dữ liệu",
        [
            "Huấn luyện DETR ĐẦY ĐỦ: nhiều epoch hơn, batch lớn hơn; thử biến thể Deformable DETR (hội tụ nhanh hơn)",
            "Xử lý mất cân bằng lớp: nhân thêm mẫu lớp hiếm, dùng loss có trọng số, thu thập thêm ảnh",
            "Augmentation an toàn (giữ đúng hướng biển): ghép ảnh, đổi sáng/tương phản, mô phỏng mưa & ban đêm",
            "Thí nghiệm DATA-EFFICIENCY: huấn luyện với 25% / 50% / 100% dữ liệu → vẽ đường cong hiệu năng",
            "Phân tích lỗi sâu: theo từng lớp, theo kích thước biển, theo điều kiện ảnh",
        ],
        "Kế hoạch cuối kỳ gồm hai nhóm việc. Nhóm thứ nhất tập trung vào mô hình và dữ liệu. Cụ thể, huấn "
        "luyện DETR đầy đủ hơn và thử biến thể Deformable DETR vốn hội tụ nhanh hơn. Xử lý mất cân bằng lớp "
        "bằng cách nhân thêm mẫu hiếm hoặc dùng loss có trọng số. Tăng cường dữ liệu một cách an toàn, giữ "
        "đúng hướng biển. Đặc biệt, làm thí nghiệm data-efficiency, tức huấn luyện với các tỉ lệ dữ liệu "
        "khác nhau để trả lời trực tiếp câu hỏi nghiên cứu ban đầu về tiết kiệm dữ liệu. Và phân tích lỗi "
        "sâu theo từng lớp và kích thước.",
        ORANGE,
        "Phần 3 — Tu",
    )

    d.bullet_slide(
        "Hướng phát triển (Đồ án cuối kỳ) — 2/2: Đánh giá & Sản phẩm",
        [
            "So sánh CÔNG BẰNG: cấp cùng ngân sách huấn luyện cho YOLO và DETR rồi đánh giá lại",
            "Kiểm tra robustness: tạo bộ test mô phỏng mưa/đêm/mờ/che + đo mức SUY GIẢM mAP",
            "Demo tương tác: ứng dụng Gradio cho phép tải ảnh lên và xem mô hình dự đoán trực tiếp",
            "Triển khai online: đẩy mô hình lên Hugging Face Hub + Spaces để ai cũng dùng thử qua trình duyệt",
            "Hoàn thiện báo cáo & notebook chạy trọn vẹn 'một cú nhấp' (tái lập hoàn toàn)",
        ],
        "Nhóm việc thứ hai hướng tới đánh giá và sản phẩm. Trước hết là so sánh công bằng bằng cách cấp cùng "
        "ngân sách huấn luyện cho cả hai mô hình. Tiếp đó là kiểm tra độ bền trên các bộ test mô phỏng điều "
        "kiện khó và đo mức suy giảm chất lượng. Về sản phẩm, nhóm dự định làm một demo tương tác bằng "
        "Gradio để tải ảnh lên và xem dự đoán ngay, rồi triển khai mô hình online trên Hugging Face để bất "
        "kỳ ai cũng dùng thử qua trình duyệt. Cuối cùng là hoàn thiện báo cáo và đảm bảo notebook chạy trọn "
        "vẹn chỉ với một cú nhấp, đúng tinh thần tái lập.",
        ORANGE,
        "Phần 3 — Tu",
    )

    d.bullet_slide(
        "Kết luận",
        [
            "Đã dựng pipeline TÁI LẬP và huấn luyện baseline cho cả YOLO và DETR trên cùng bộ dữ liệu",
            "YOLOv8n: mAP@0.5 = 0.954 · 113 FPS · 6.25 MB → rất phù hợp triển khai thực tế trên xe",
            "DETR baseline còn thấp do điều kiện huấn luyện hạn chế — KHÔNG phải do bản chất mô hình",
            "Trả lời được một phần câu hỏi: hiện YOLO 'robust + data-efficient' hơn trong điều kiện hạn chế",
            "Cuối kỳ: so sánh công bằng · tăng robustness · data-efficiency · demo Gradio · deploy Hugging Face",
            "Cảm ơn thầy/cô và các bạn đã lắng nghe — nhóm xin sẵn sàng trả lời câu hỏi!",
        ],
        "Để kết luận, ở giai đoạn giữa kỳ nhóm đã dựng được một pipeline tái lập và huấn luyện baseline cho "
        "cả hai mô hình trên cùng bộ dữ liệu. YOLO cho kết quả rất mạnh và phù hợp triển khai thực tế trên "
        "xe. DETR còn thấp nhưng là do điều kiện huấn luyện hạn chế chứ không phải bản chất. Như vậy nhóm đã "
        "trả lời được một phần câu hỏi nghiên cứu: trong điều kiện hạn chế, YOLO vừa chính xác vừa tiết kiệm "
        "hơn. Trong đồ án cuối kỳ, nhóm sẽ so sánh công bằng, tăng độ bền, làm thí nghiệm data-efficiency, "
        "và xây demo cùng triển khai online. Nhóm em xin cảm ơn thầy cô và các bạn, và rất sẵn lòng trả lời câu hỏi.",
        ORANGE,
        "Phần 3 — Tu",
    )

    d.save()
    n = len(d.prs.slides._sldIdLst)
    print(f"✔ Đã tạo {n} slide → {OUT_PATH}")


if __name__ == "__main__":
    build()
